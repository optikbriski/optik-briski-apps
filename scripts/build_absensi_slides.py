#!/usr/bin/env python3
"""Buat presentasi Alur Absensi (bahasa awam) → .pptx untuk diunggah ke Google Slides."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parents[1] / "docs" / "Alur-Absensi-Optik-B-Riski.pptx"

NAVY = RGBColor(0x0B, 0x1F, 0x3A)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x4A, 0x55, 0x68)
SOFT = RGBColor(0xF5, 0xF7, 0xFA)
LINE = RGBColor(0xD8, 0xDE, 0xE8)


def _set_run(run, text, size=18, bold=False, color=INK, font="Arial"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def _add_bg(slide, color):
    fill = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.333), Inches(7.5)  # rectangle
    )
    fill.fill.solid()
    fill.fill.fore_color.rgb = color
    fill.line.fill.background()
    # send to back
    spTree = slide.shapes._spTree
    sp = fill._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def _bar(slide, y=0):
    shape = slide.shapes.add_shape(1, Inches(0), Inches(y), Inches(13.333), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    return shape


def _title_bar(slide, title, subtitle=None):
    _bar(slide)
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.55))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    _set_run(r, title, 26, True, WHITE)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.4))
        sp = sub.text_frame.paragraphs[0]
        sr = sp.add_run()
        _set_run(sr, subtitle, 14, False, MUTED)


def _card(slide, left, top, width, height, fill=SOFT):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1)
    return shape


def _bullets(slide, left, top, width, height, lines, size=16):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(8)
        r = p.add_run()
        _set_run(r, line, size, False, INK)
    return box


def _heading(slide, left, top, width, text, size=18, color=NAVY):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.4))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    _set_run(r, text, size, True, color)
    return box


def add_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, NAVY)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(1.2))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    _set_run(r, "Alur Absensi", 44, True, WHITE)
    sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(11.5), Inches(1.2))
    sp = sub.text_frame.paragraphs[0]
    sr = sp.add_run()
    _set_run(sr, "Optik B. Riski — panduan fitur untuk Admin & Karyawan\n(bahasa awam, sesuai sistem yang sudah jalan sekarang)", 18, False, GOLD)
    foot = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.4))
    fp = foot.text_frame.paragraphs[0]
    fr = fp.add_run()
    _set_run(fr, "Web Admin  •  Absensi Toko  •  APK Karyawan  •  Monitor Absensi", 14, False, WHITE)


def add_dua_perangkat(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _title_bar(slide, "1. Dua perangkat yang dipakai", "Absensi butuh kerja sama layar Admin toko + HP karyawan")
    _card(slide, 0.5, 1.6, 5.9, 5.2)
    _card(slide, 6.9, 1.6, 5.9, 5.2)
    _heading(slide, 0.75, 1.8, 5.4, "Layar Admin (Web / tablet toko)")
    _bullets(
        slide,
        0.75,
        2.4,
        5.4,
        4.2,
        [
            "• Menu Absensi (Absensi Toko)",
            "• Menampilkan QR Absensi yang terus berganti",
            "• Menunggu karyawan scan di HP",
            "• Untuk masuk: minta wajah (+ PIN jika ada)",
            "• Untuk pulang: tercatat otomatis setelah scan + lokasi OK",
            "• Menu Monitor Absensi: tinjau foto masuk (Valid / Mencurigakan)",
        ],
        15,
    )
    _heading(slide, 7.15, 1.8, 5.4, "HP Karyawan (APK)")
    _bullets(
        slide,
        7.15,
        2.4,
        5.4,
        4.2,
        [
            "• Menu Absensi di aplikasi karyawan",
            "• Sekali: daftarkan wajah (di area toko)",
            "• Setiap absen: Scan QR Absensi di layar Admin",
            "• GPS HP harus di dalam area toko (QR saja tidak cukup)",
            "• Setelah masuk: lokasi HP dipantau selama shift",
            "• Pulang: scan QR lagi di area toko (tanpa wajah)",
        ],
        15,
    )


def add_menu_web(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _title_bar(slide, "2. Menu di Web Admin yang terkait absensi", "Ini yang Admin lihat / pakai sehari-hari")
    items = [
        ("Absensi", "Layar utama toko: tampilkan QR, tangkap absen masuk/pulang karyawan."),
        ("QR Absensi Toko", "Cadangan menampilkan QR toko saja (jika perlu layar QR terpisah)."),
        ("Monitor Absensi", "Antrian foto absen masuk — Admin putuskan Valid atau Mencurigakan."),
        ("Geofence Toko", "Gambar area toko di peta. Absensi hanya lolos jika GPS HP di dalam area ini."),
        ("Jadwal Kerja", "Jam masuk / pulang / libur per karyawan. Tanpa jadwal hari itu, absen masuk ditolak."),
        ("Data Karyawan", "Toko karyawan, status Aktif, PIN absensi, foto wajah terdaftar."),
    ]
    y = 1.5
    for title, desc in items:
        _card(slide, 0.5, y, 12.3, 0.85)
        _heading(slide, 0.7, y + 0.08, 3.2, title, 16, GOLD if False else NAVY)
        box = slide.shapes.add_textbox(Inches(3.8), Inches(y + 0.18), Inches(8.7), Inches(0.55))
        p = box.text_frame.paragraphs[0]
        r = p.add_run()
        _set_run(r, desc, 14, False, INK)
        y += 0.95


def add_dari_mana_qr(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _title_bar(
        slide,
        "3. Dari mana QR Absensi muncul?",
        "Bukan di-print. QR hidup di layar Admin dan berganti sendiri",
    )
    steps = [
        ("1", "Admin buka menu Absensi", "Sistem mengenali toko perangkat Admin (cabang / pusat)."),
        ("2", "Layar menampilkan QR Absensi Toko", "QR terikat ke toko itu. Karyawan toko lain tidak bisa pakai."),
        ("3", "QR berganti otomatis tiap beberapa detik", "Screenshot lama tidak bisa dipakai. Selalu scan layar yang sedang tampil."),
        ("4", "Karyawan scan dari HP", "Kalau GPS di dalam area toko → Admin langsung “menangkap” karyawan itu."),
        ("5", "Setelah absen selesai", "Layar Admin kembali ke QR, siap karyawan berikutnya."),
    ]
    y = 1.45
    for num, title, desc in steps:
        circ = slide.shapes.add_shape(3, Inches(0.55), Inches(y + 0.1), Inches(0.55), Inches(0.55))
        circ.fill.solid()
        circ.fill.fore_color.rgb = NAVY
        circ.line.fill.background()
        t = circ.text_frame
        t.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = t.paragraphs[0].add_run()
        _set_run(r, num, 16, True, WHITE)
        _heading(slide, 1.35, y, 11, title, 17)
        box = slide.shapes.add_textbox(Inches(1.35), Inches(y + 0.38), Inches(11), Inches(0.4))
        p = box.text_frame.paragraphs[0]
        rr = p.add_run()
        _set_run(rr, desc, 14, False, MUTED)
        y += 1.05


def add_syarat(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _title_bar(slide, "4. Syarat sebelum absen masuk berhasil", "Semua harus terpenuhi — kalau satu gagal, absen ditolak")
    cols = [
        (
            "Akun & toko",
            [
                "• Status karyawan: Aktif",
                "• Toko karyawan sudah terisi",
                "• Toko sama dengan QR yang discan",
            ],
        ),
        (
            "Jadwal hari ini",
            [
                "• Ada jadwal kerja hari itu",
                "• Bukan hari libur",
                "• Jam masuk sudah diisi di jadwal",
            ],
        ),
        (
            "Lokasi & identitas",
            [
                "• Scan QR di layar Admin toko",
                "• GPS HP di dalam area toko",
                "• Wajah sudah terdaftar",
                "• PIN benar (jika karyawan punya PIN)",
            ],
        ),
    ]
    x = 0.5
    for title, lines in cols:
        _card(slide, x, 1.6, 4.0, 5.0)
        _heading(slide, x + 0.25, 1.85, 3.5, title, 18)
        _bullets(slide, x + 0.25, 2.5, 3.5, 3.8, lines, 15)
        x += 4.2
    note = slide.shapes.add_textbox(Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4))
    p = note.text_frame.paragraphs[0]
    r = p.add_run()
    _set_run(r, "Catatan: komputer Admin tidak pakai GPS. Bukti lokasi selalu dari GPS HP karyawan saat scan.", 13, False, MUTED)


def add_alur_masuk(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _title_bar(slide, "5. Alur Absen Masuk (langkah demi langkah)", "Dari QR di web → sampai shift aktif di HP")
    left = [
        "Di Admin (Web / tablet)",
        "1. Buka menu Absensi",
        "2. Biarkan QR tampil di layar",
        "3. Tunggu notifikasi: karyawan sudah scan + lokasi OK",
        "4. Kalau wajah belum terdaftar → daftarkan dulu",
        "5. Kalau ada PIN → karyawan isi PIN",
        "6. Ambil foto wajah (kamera Admin)",
        "7. Sistem catat Absen masuk",
        "8. Foto masuk masuk antrian Monitor Absensi",
        "9. Layar kembali ke QR",
    ]
    right = [
        "Di HP Karyawan",
        "1. Login (status Aktif)",
        "2. Buka Absensi",
        "3. Pastikan sudah di dalam toko",
        "4. Ketuk Scan QR Absensi (masuk/pulang)",
        "5. Arahkan kamera ke QR di layar Admin",
        "6. Jika lokasi OK → muncul pesan lanjut ke Admin",
        "7. Selesaikan wajah (+ PIN) di perangkat Admin",
        "8. HP mendeteksi masuk tercatat",
        "9. Pantauan lokasi shift dimulai",
    ]
    _card(slide, 0.4, 1.45, 6.1, 5.6)
    _card(slide, 6.8, 1.45, 6.1, 5.6)
    _heading(slide, 0.65, 1.6, 5.6, left[0], 17, NAVY)
    _bullets(slide, 0.65, 2.15, 5.6, 4.7, left[1:], 14)
    _heading(slide, 7.05, 1.6, 5.6, right[0], 17, NAVY)
    _bullets(slide, 7.05, 2.15, 5.6, 4.7, right[1:], 14)


def add_alur_pulang(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _title_bar(slide, "6. Alur Absen Pulang", "Lebih sederhana: tidak perlu wajah di Admin")
    _card(slide, 0.5, 1.5, 12.3, 5.4)
    lines = [
        "1. Karyawan sudah punya shift aktif (sudah absen masuk hari itu).",
        "2. Admin tetap menampilkan QR Absensi di menu Absensi.",
        "3. Karyawan di HP: Scan QR Absensi lagi — harus masih di area toko (GPS).",
        "4. Admin otomatis mencatat Absen pulang (tanpa PIN, tanpa foto wajah).",
        "5. Shift ditutup. Di HP, pantauan lokasi berhenti.",
        "",
        "Yang sering ditolak saat pulang:",
        "• Belum waktunya pulang (mengikuti jam pulang di Jadwal Kerja)",
        "• GPS di luar area toko meski QR berhasil terbaca",
        "• Sudah absen pulang hari ini (hanya 1× per tanggal)",
    ]
    _bullets(slide, 0.85, 1.8, 11.6, 4.8, lines, 16)


def add_daftar_wajah(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _title_bar(slide, "7. Daftar wajah — dilakukan sekali", "Bisa dari HP karyawan atau dari Absensi Admin")
    _card(slide, 0.5, 1.55, 6.0, 5.2)
    _card(slide, 6.8, 1.55, 6.0, 5.2)
    _heading(slide, 0.75, 1.8, 5.5, "Dari HP Karyawan")
    _bullets(
        slide,
        0.75,
        2.4,
        5.5,
        4.0,
        [
            "• Buka Absensi → Daftarkan wajah",
            "• Harus berada di area toko",
            "• Ikuti petunjuk kamera (liveness)",
            "• Setelah sukses: status “Wajah terdaftar”",
            "• Bisa daftar ulang jika foto kurang jelas",
        ],
        15,
    )
    _heading(slide, 7.05, 1.8, 5.5, "Dari Admin Absensi")
    _bullets(
        slide,
        7.05,
        2.4,
        5.5,
        4.0,
        [
            "• Karyawan scan QR dulu",
            "• Jika wajah belum ada, Admin diminta daftarkan",
            "• PIN diminta jika karyawan punya PIN",
            "• Setelah daftar, minta scan QR lagi untuk absen masuk",
            "• Foto terdaftar dipakai banding di Monitor Absensi",
        ],
        15,
    )


def add_monitor(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _title_bar(slide, "8. Monitor Absensi (setelah masuk tercatat)", "Keputusan Valid / Mencurigakan ada di sini — bukan di layar QR")
    _card(slide, 0.5, 1.55, 12.3, 5.3)
    lines = [
        "Setelah Absen masuk berhasil di Absensi Toko:",
        "• Foto wajah masuk masuk antrian Monitor Absensi",
        "• Admin membuka menu Monitor Absensi",
        "• Bandingkan foto masuk dengan foto wajah terdaftar",
        "",
        "Pilihan Admin:",
        "• Valid / Aman → poin kehadiran diproses sesuai aturan",
        "• Mencurigakan / Curang → sanksi poin & peringatan sesuai aturan perusahaan",
        "",
        "Poin tidak langsung “final” saat absen — menunggu tinjauan Monitor Absensi.",
    ]
    _bullets(slide, 0.85, 1.85, 11.6, 4.8, lines, 16)


def add_pantauan(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _title_bar(slide, "9. Pantauan lokasi selama shift", "Berjalan di HP karyawan setelah absen masuk")
    _card(slide, 0.5, 1.5, 12.3, 5.4)
    lines = [
        "• Setelah masuk tercatat, HP memantau apakah karyawan masih di area toko.",
        "• Idealnya izin lokasi “Selalu”, notifikasi aktif, dan (Android) optimasi baterai dimatikan untuk app.",
        "• Keluar area toko tanpa ijin/cuti disetujui → muncul peringatan.",
        "• Ada ijin/cuti disetujui hari itu → pantauan keluar area tidak memicu peringatan.",
        "• Sebelum jam standby shift (pagi / siang sesuai aturan), peringatan keluar belum aktif.",
        "• Absen pulang → pantauan otomatis berhenti.",
        "",
        "Penting: force-stop aplikasi atau cabut izin lokasi bisa menghentikan pantauan.",
    ]
    _bullets(slide, 0.85, 1.8, 11.6, 4.8, lines, 16)


def add_ringkas(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _title_bar(slide, "10. Ringkasan satu layar", "Cara ingat alur absensi sehari-hari")
    rows = [
        ("Siapkan", "Admin buka Absensi → QR tampil. Karyawan sudah Aktif, punya jadwal, wajah terdaftar."),
        ("Masuk", "HP scan QR di area toko → Admin wajah (+ PIN) → masuk tercatat → pantauan mulai."),
        ("Kerja", "Lokasi HP dipantau. Monitor Absensi meninjau foto masuk (Valid / Mencurigakan)."),
        ("Pulang", "HP scan QR lagi di area toko → Admin catat pulang otomatis → pantauan berhenti."),
    ]
    y = 1.5
    for title, desc in rows:
        _card(slide, 0.5, y, 12.3, 1.2)
        _heading(slide, 0.75, y + 0.15, 2.2, title, 18, NAVY)
        box = slide.shapes.add_textbox(Inches(3.0), Inches(y + 0.3), Inches(9.4), Inches(0.7))
        p = box.text_frame.paragraphs[0]
        r = p.add_run()
        _set_run(r, desc, 15, False, INK)
        y += 1.35


def add_faq(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _title_bar(slide, "11. Pertanyaan umum", "Jawaban singkat untuk Admin & Karyawan")
    faqs = [
        ("Kenapa QR selalu berubah?", "Supaya screenshot lama / foto QR kemarin tidak bisa dipakai absen."),
        ("Kenapa QR terbaca tapi tetap ditolak?", "Karena GPS HP di luar area toko. Harus masuk area dulu, lalu scan ulang."),
        ("Kenapa Mac/PC Admin tidak cek GPS?", "Komputer biasanya tidak punya GPS. Bukti lokasi diambil dari HP karyawan."),
        ("Di mana keputusan Valid / Curang?", "Di menu Monitor Absensi — bukan di layar QR Absensi."),
        ("Pulang perlu wajah?", "Tidak. Cukup scan QR + GPS di area toko (dan sudah waktunya pulang)."),
        ("Siapa yang menampilkan QR?", "Layar Admin di toko (menu Absensi). Karyawan hanya scan dari HP."),
    ]
    y = 1.4
    for q, a in faqs:
        _heading(slide, 0.55, y, 12.2, q, 15, NAVY)
        box = slide.shapes.add_textbox(Inches(0.55), Inches(y + 0.32), Inches(12.2), Inches(0.35))
        p = box.text_frame.paragraphs[0]
        r = p.add_run()
        _set_run(r, a, 13, False, MUTED)
        y += 0.9


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_cover(prs)
    add_dua_perangkat(prs)
    add_menu_web(prs)
    add_dari_mana_qr(prs)
    add_syarat(prs)
    add_alur_masuk(prs)
    add_alur_pulang(prs)
    add_daftar_wajah(prs)
    add_monitor(prs)
    add_pantauan(prs)
    add_ringkas(prs)
    add_faq(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
